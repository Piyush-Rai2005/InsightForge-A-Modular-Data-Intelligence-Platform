"""
Generate a PowerPoint deck from analysis results.
"""
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def generate_pptx(report: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── Title slide ─────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "InsightForge Analysis Report"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x63, 0xD3, 0x96)

    # ── Executive Summary slide ─────────────────────────────────────────────
    exec_sum = report.get("exec_summary", "")
    if exec_sum:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = "Executive Summary"
        tp.font.size = Pt(28)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(0x63, 0xD3, 0x96)

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(5))
        bf = body_box.text_frame
        bf.word_wrap = True
        for line in exec_sum.split("\n"):
            p = bf.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(14)
            p.space_after = Pt(8)

    # ── Model Comparison slide ──────────────────────────────────────────────
    models = report.get("model_comparison", [])
    if models:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = "Model Performance"
        tp.font.size = Pt(28)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(0x63, 0xD3, 0x96)

        rows = len(models) + 1
        table_shape = slide.shapes.add_table(rows, 2, Inches(2), Inches(2), Inches(9), Inches(0.5 * rows))
        table = table_shape.table
        table.cell(0, 0).text = "Model"
        table.cell(0, 1).text = "Accuracy"
        for i, m in enumerate(models):
            table.cell(i + 1, 0).text = m["model"]
            try:
                table.cell(i + 1, 1).text = f"{float(m['accuracy']) * 100:.2f}%"
            except (ValueError, TypeError):
                table.cell(i + 1, 1).text = str(m.get('accuracy', 'N/A'))

    # ── Recommendations slide ───────────────────────────────────────────────
    recs = report.get("recommendations", [])
    if recs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tp = title_box.text_frame.paragraphs[0]
        tp.text = "Recommendations"
        tp.font.size = Pt(28)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(0x63, 0xD3, 0x96)

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(5))
        bf = body_box.text_frame
        bf.word_wrap = True
        for rec in recs:
            p = bf.add_paragraph()
            p.text = f"• {rec}"
            p.font.size = Pt(14)
            p.space_after = Pt(10)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
